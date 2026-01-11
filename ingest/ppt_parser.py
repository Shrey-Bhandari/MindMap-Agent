# convert PPTs → normalized, context-aware slide records.
from pptx import Presentation
from pathlib import Path
from typing import List, Dict, Optional


def extract_slide_text(slide) -> List[str]:
    """
    Extract visible text from a slide.
    Images, charts, and diagrams are ignored by design.
    """
    texts = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                texts.append(text)

    return texts


def load_ppt(ppt_path: Path) -> List[Dict]:
    """
    Load a PPT file and extract raw slide content.
    """
    prs = Presentation(ppt_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        slide_record = {
            "slide_number": idx + 1,
            "title": slide.shapes.title.text.strip()
            if slide.shapes.title and slide.shapes.title.text
            else None,
            "raw_text": extract_slide_text(slide),
            "has_image": any(
                shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
                for shape in slide.shapes
            ),
        }
        slides_data.append(slide_record)

    return slides_data
